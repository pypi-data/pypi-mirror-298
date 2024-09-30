import os
import json
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ..abstract_handler import AbstractHandler

class MicrosoftPowerPointReaderHandler(AbstractHandler):
    
    def handle(self, request: dict) -> dict:
        print("Processing PPTX file...")

        # Initialize variables
        text_content = ''
        media_files = []

        # Load the presentation from the path specified in the request
        presentation_path = request.get("path", None)
        presentation = Presentation(presentation_path)

        # Determine the write file path and output folder
        write_file_path = request.get("write_file_path", None)
        if write_file_path:
            base_folder = os.path.dirname(write_file_path)
            file_id = os.path.basename(write_file_path).split('.')[0]
            output_folder = os.path.join(base_folder, file_id)
            media_folder = os.path.join(output_folder, "media")
            os.makedirs(media_folder, exist_ok=True)
        
        # Initialize the metadata dictionary
        if "metadata" not in request:
            request["metadata"] = {}

        # Iterate through each slide in the presentation
        for slide_number, slide in enumerate(presentation.slides, start=1):
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content += shape.text + '\n'

                if request.get("extract_media", False):
                    # Extract images and videos
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        image_file_name = f"slide_{slide_number}_image_{shape.shape_id}.{image.ext}"
                        image_file_path = os.path.join(media_folder, image_file_name)
                        with open(image_file_path, "wb") as f:
                            f.write(image.blob)
                        media_files.append({"type": "image", "path": image_file_path, "slide_number": slide_number, "shape_id": shape.shape_id})

                    elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                        media = shape
                        media_file_name = f"slide_{slide_number}_media_{shape.shape_id}"

                        if media.media_format == 'video':
                            media_file_path = os.path.join(media_folder, media_file_name + ".mp4")  # Assuming video format is mp4
                        elif media.media_format == 'audio':
                            media_file_path = os.path.join(media_folder, media_file_name + ".mp3")  # Assuming audio format is mp3
                        else:
                            media_file_path = os.path.join(media_folder, media_file_name)

                        with open(media_file_path, "wb") as f:
                            f.write(media.blob)
                        media_files.append({"type": media.media_format, "path": media_file_path, "slide_number": slide_number, "shape_id": shape.shape_id})

            # Extract text from speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                text_content += notes_slide.notes_text_frame.text + '\n'

        # Save the transcript to a text file
        transcript_file_path = os.path.join(output_folder, 'transcript.txt')
        with open(transcript_file_path, 'w') as transcript_file:
            transcript_file.write(text_content)

        # Create metadata with traceability for graph database
        metadata = {
            "original_file": presentation_path,
            "transcript_file": transcript_file_path,
            "media_files": media_files
        }
        
        # Save metadata to a json file
        metadata_file_path = os.path.join(output_folder, 'metadata.json')
        with open(metadata_file_path, 'w') as metadata_file:
            json.dump(metadata, metadata_file, indent=4)

        # Update the request metadata with the new metadata
        request["metadata"][presentation_path] = metadata

        # Update the request with the extracted text
        request.update({"text": text_content})

        # Call the next handler in the chain
        return super().handle(request)
